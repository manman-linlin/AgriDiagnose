"""AgriDiagnose Web 前端答辩 PPT 生成脚本"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from lxml import etree
import os

# ═══════════════════════════════════════════════════════
# 0. 初始化 — 用参考模板作为母版基底
# ═══════════════════════════════════════════════════════
TEMPLATE_PATH = 'D:/个人文件/AgriDiagnose/AgriDiagnose/_clean.pptx'
SAVE_PATH = 'D:/个人文件/AgriDiagnose/AgriDiagnose/AgriDiagnose-Web前端答辩.pptx'

prs = Presentation(TEMPLATE_PATH)

# 布局
LO = prs.slide_masters[0].slide_layouts
COVER_LAYOUT = LO[0]          # 封面
TITLE_CONTENT_LAYOUT = LO[3]  # 标题和内容

# ═══════════════════════════════════════════════════════
# 常量
# ═══════════════════════════════════════════════════════
SLIDE_W = 12192000  # 13.33 in
SLIDE_H = 6858000   # 7.5 in
MARGIN = Inches(0.4)
CONTENT_W = SLIDE_W - MARGIN * 2

# 颜色
GOLD    = RGBColor(0xFF, 0xC0, 0x00)
BLUE    = RGBColor(0x5B, 0x9B, 0xD5)
DARK    = RGBColor(0x44, 0x54, 0x6A)
GRAY    = RGBColor(0x59, 0x59, 0x59)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x00, 0x00, 0x00)
LGRAY   = RGBColor(0xF2, 0xF4, 0xF7)
MLGRAY  = RGBColor(0xE0, 0xE0, 0xE0)
PLACEHOLDER_BG = RGBColor(0xF8, 0xF9, 0xFA)
PLACEHOLDER_BD = RGBColor(0xCC, 0xCC, 0xCC)
PLACEHOLDER_TXT = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT_ORANGE = RGBColor(0xED, 0x7D, 0x31)
ACCENT_GREEN  = RGBColor(0x70, 0xAD, 0x47)

# 字体
TITLE_FONT = '思源宋体 CN Heavy'  # 模板标题字体
BODY_FONT  = '微软雅黑'

# ═══════════════════════════════════════════════════════
# 辅助函数
# ═══════════════════════════════════════════════════════

def add_tb(slide, l, t, w, h, text, font=BODY_FONT, size=Pt(16), bold=False,
           color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """添加文本框"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = align
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = size
    p.font.bold = bold
    try:
        p.font.color.rgb = color
    except:
        pass
    p.alignment = align
    return tb

def add_multiline(slide, l, t, w, h, lines, font=BODY_FONT, align=PP_ALIGN.LEFT):
    """添加多行文本，lines=[(text, size, bold, color), ...]"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        text, size, bold, color = ln
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = font
        p.font.size = size
        p.font.bold = bold
        try:
            p.font.color.rgb = color
        except:
            pass
        p.alignment = align
        p.space_after = Pt(3)
    return tb

def add_card(slide, l, t, w, h, num_text, desc_text):
    """数字指标卡片"""
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    rect.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    rect.line.width = Pt(0.5)
    # 数字
    add_tb(slide, l, t + h//2 - Inches(0.25), w, Inches(0.3),
           num_text, size=Pt(24), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # 描述
    add_tb(slide, l, t + h//2 + Inches(0.08), w, Inches(0.2),
           desc_text, size=Pt(13), color=GRAY, align=PP_ALIGN.CENTER)

def add_placeholder(slide, l, t, w, h, label="截图"):
    """截图占位符"""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = PLACEHOLDER_BG
    rect.line.color.rgb = PLACEHOLDER_BD
    rect.line.width = Pt(1)
    # 虚线
    rect.line.dash_style = 2
    tf = rect.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[ {label} ]"
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.color.rgb = PLACEHOLDER_TXT
    p.alignment = PP_ALIGN.CENTER
    return rect

def add_page_title(slide, num, title):
    """页面标题"""
    add_tb(slide, MARGIN, Inches(0.2), Inches(5), Inches(0.55),
           f"{num:02d}  {title}", font=TITLE_FONT, size=Pt(34), color=DARK)

def add_bottom_bar(slide, text):
    """底部信息条"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.45),
                                  SLIDE_W, Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    bar.line.fill.background()
    add_tb(slide, MARGIN, SLIDE_H - Inches(0.4), SLIDE_W - MARGIN*2, Inches(0.35),
           text, size=Pt(11), color=GRAY, align=PP_ALIGN.CENTER)

def add_annotation(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT):
    """添加注释文字（小字号关键词列表）"""
    return add_multiline(slide, l, t, w, h, lines, font=BODY_FONT, align=align)

# ═══════════════════════════════════════════════════════
# 第 1 页 — 工作概览
# ═══════════════════════════════════════════════════════
print("生成第 1 页...")
s1 = prs.slides.add_slide(COVER_LAYOUT)

# 标题行
add_tb(s1, MARGIN, Inches(0.5), CONTENT_W, Inches(0.7),
       "Web 前端工作汇报 — 农作物病虫害智能诊断系统",
       font=TITLE_FONT, size=Pt(36), color=DARK, align=PP_ALIGN.CENTER)

# 2×3 卡片 + 中央截图
CARD_W = Inches(2.6)
CARD_H = Inches(0.9)
CARD_GAP = Inches(0.3)
TOP_CARDS_Y = Inches(1.5)
MID_Y = Inches(2.7)
BOT_CARDS_Y = Inches(5.1)
CENTER_IMG_W = Inches(5.5)
CENTER_IMG_H = Inches(2.1)

# 上排三个卡片
cards_top = ["55 个", "7,700 行", "13 个"]
cards_top_desc = ["源文件", "代码", "页面"]
total_top_w = 3 * CARD_W + 2 * CARD_GAP
start_x = (SLIDE_W - total_top_w) // 2
for i in range(3):
    add_card(s1, start_x + i*(CARD_W + CARD_GAP), TOP_CARDS_Y, CARD_W, CARD_H,
             cards_top[i], cards_top_desc[i])

# 中央目录截图
center_x = (SLIDE_W - CENTER_IMG_W) // 2
add_placeholder(s1, center_x, MID_Y, CENTER_IMG_W, CENTER_IMG_H,
                "项目目录结构截图 (VSCode 文件树)")

# 下排三个卡片
cards_bot = ["6 大", "9 个", "147 个"]
cards_bot_desc = ["功能模块", "通用组件", "设计 Token"]
for i in range(3):
    add_card(s1, start_x + i*(CARD_W + CARD_GAP), BOT_CARDS_Y, CARD_W, CARD_H,
             cards_bot[i], cards_bot_desc[i])

# 底部技术栈
add_tb(s1, MARGIN, SLIDE_H - Inches(0.8), CONTENT_W, Inches(0.4),
       "Vue 3 + Pinia + ECharts  ·  零构建  ·  ESM importmap  ·  CDN 直载",
       size=Pt(13), color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# 第 2 页 — 架构设计
# ═══════════════════════════════════════════════════════
print("生成第 2 页...")
s2 = prs.slides.add_slide(TITLE_CONTENT_LAYOUT)
add_page_title(s2, 2, "架构设计")

# 左侧架构图区域
ARCH_LEFT = MARGIN
ARCH_TOP = Inches(1.2)
ARCH_W = Inches(6.5)
ARCH_H = Inches(5.0)

# 五层架构用圆角矩形 + 箭头
BOX_W = Inches(5.2)
BOX_H = Inches(0.65)
BOX_GAP = Inches(0.18)
ARR_H = Inches(0.2)

layer_x = ARCH_LEFT + Inches(0.6)
layers = [
    ("路由层 · 状态层 · 网络层", "Hash 路由  |  6 个 Pinia Store  |  Bearer Token 自动注入"),
    ("布局组件层 · 通用组件层", "Topbar / Sidebar / BottomNav  |  9 个组件 + 2 个 Composable"),
    ("功能页面层", "6 大功能模块 · 13 个可交互页面"),
]

y = ARCH_TOP + Inches(0.5)
for i, (title, desc) in enumerate(layers):
    # 层背景
    rect = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                layer_x, y, BOX_W, BOX_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xEE, 0xF1, 0xF6)
    rect.line.color.rgb = BLUE
    rect.line.width = Pt(1)
    # 标题
    add_tb(s2, layer_x + Inches(0.15), y + Inches(0.05), BOX_W - Inches(0.3), Inches(0.3),
           title, size=Pt(15), bold=True, color=DARK)
    # 描述
    add_tb(s2, layer_x + Inches(0.15), y + Inches(0.32), BOX_W - Inches(0.3), Inches(0.28),
           desc, size=Pt(11), color=GRAY)
    # 层间箭头
    if i < len(layers) - 1:
        arrow_y = y + BOX_H
        arrow = s2.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     layer_x + BOX_W//2 - Inches(0.15),
                                     arrow_y + Inches(0.02),
                                     Inches(0.3), ARR_H)
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = BLUE
        arrow.line.fill.background()
    y += BOX_H + ARR_H + BOX_GAP

# 底部标注：数据流方向
add_tb(s2, layer_x, y + Inches(0.1), BOX_W, Inches(0.3),
       "数据流：用户操作 → 状态管理 → API 请求 → 后端响应 → 视图更新",
       size=Pt(11), color=GRAY, align=PP_ALIGN.CENTER)

# 右侧三点说明
RIGHT_X = Inches(7.8)
RIGHT_Y = Inches(1.3)
RIGHT_W = Inches(4.8)

right_items = [
    ("📱  响应式策略", "桌面端 / 平板端 / 移动端\n三级断点自适应，底部导航自动切换"),
    ("🎨  设计 Token 体系", "147 个 CSS 变量 → 5 级分层\nvariables → layout → components → animations → responsive"),
    ("🏗  构建方案", "零编译 · ESM importmap · CDN 直载\nVue 3 + Pinia + ECharts"),
]

item_y = RIGHT_Y
for title, desc in right_items:
    # 卡片背景
    rect = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                RIGHT_X, item_y, RIGHT_W, Inches(1.3))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    rect.line.color.rgb = MLGRAY
    rect.line.width = Pt(0.5)

    add_tb(s2, RIGHT_X + Inches(0.2), item_y + Inches(0.15), RIGHT_W - Inches(0.4), Inches(0.35),
           title, size=Pt(18), bold=True, color=DARK)
    add_tb(s2, RIGHT_X + Inches(0.2), item_y + Inches(0.55), RIGHT_W - Inches(0.4), Inches(0.7),
           desc, size=Pt(13), color=GRAY)
    item_y += Inches(1.5)

add_bottom_bar(s2, "技术选型：Vue 3 + Pinia + ECharts  |  全部后端 API 已完成联调")

# ═══════════════════════════════════════════════════════
# 第 3 页 — 核心功能 · 用户端
# ═══════════════════════════════════════════════════════
print("生成第 3 页...")
s3 = prs.slides.add_slide(TITLE_CONTENT_LAYOUT)
add_page_title(s3, 3, "核心功能 — 用户端")

IMG_W = Inches(3.0)
IMG_H = Inches(1.65)
IMG_GAP = Inches(0.5)

# 2-1-2 菱形对称布局
ROW1_Y = Inches(1.05)
ROW2_Y = Inches(2.95)
ROW3_Y = Inches(4.85)

# 上行两个
for i in range(2):
    x = (SLIDE_W - 2*IMG_W - IMG_GAP) // 2 + i*(IMG_W + IMG_GAP)
    add_placeholder(s3, x, ROW1_Y, IMG_W, IMG_H, f"用户端截图 {i+1}")

# 中行一个（居中）
add_placeholder(s3, (SLIDE_W - IMG_W)//2, ROW2_Y, IMG_W, IMG_H, "用户端截图 3")

# 下行两个
for i in range(2):
    x = (SLIDE_W - 2*IMG_W - IMG_GAP) // 2 + i*(IMG_W + IMG_GAP)
    add_placeholder(s3, x, ROW3_Y, IMG_W, IMG_H, f"用户端截图 {i+4}")

# 注释（外侧关键词）—— 每个截图右侧标注模块名 + 功能关键词
annot_data = [
    # (row_y, x_index, module_name, keywords)
    (ROW1_Y, 0, "智能诊断", "拖拽上传\nTop-3 结果\n置信度可视化"),
    (ROW1_Y, 1, "AI 对话",   "SSE 流式输出\n多模态复核\n防治建议卡片"),
    (ROW3_Y, 0, "数据贡献",  "扩展已有类别\n新增病害类型\n标签上传审核"),
    (ROW3_Y, 1, "历史统计",  "诊断记录管理\nECharts 饼图\nCSV 数据导出"),
]
for ry, xi, mname, kws in annot_data:
    x = (SLIDE_W - 2*IMG_W - IMG_GAP) // 2 + xi*(IMG_W + IMG_GAP)
    add_tb(s3, x + IMG_W + Inches(0.1), ry + Inches(0.15), Inches(1.5), Inches(0.3),
           mname, size=Pt(11), bold=True, color=BLUE)
    add_tb(s3, x + IMG_W + Inches(0.1), ry + Inches(0.5), Inches(1.5), Inches(0.7),
           kws, size=Pt(10), color=GRAY)

# 中行注释（左右两侧）
add_tb(s3, (SLIDE_W - IMG_W)//2 - Inches(1.6), ROW2_Y + Inches(0.3), Inches(1.5), Inches(0.8),
       "38 类病害图鉴\n作物/分类/关键词\n三维筛选 + 详情抽屉",
       size=Pt(10), color=GRAY, align=PP_ALIGN.RIGHT)
add_tb(s3, (SLIDE_W + IMG_W)//2 + Inches(0.1), ROW2_Y + Inches(0.3), Inches(1.4), Inches(0.8),
       "病害百科",
       size=Pt(11), bold=True, color=BLUE)

add_bottom_bar(s3, "诊断 → 对话 → 百科查询 → 数据反馈，覆盖完整业务闭环")

# ═══════════════════════════════════════════════════════
# 第 4 页 — 核心功能 · 管理端
# ═══════════════════════════════════════════════════════
print("生成第 4 页...")
s4 = prs.slides.add_slide(TITLE_CONTENT_LAYOUT)
add_page_title(s4, 4, "核心功能 — 管理端")

S_IMG_W = Inches(2.7)
S_IMG_H = Inches(1.45)
S_IMG_GAP = Inches(0.25)

SR1_Y = Inches(1.05)
SR2_Y = Inches(2.8)
SR3_Y = Inches(4.55)
SR4_Y = Inches(6.25)  # 末行（系统日志单独一行居中）

# 上行三个
srow1_labels = ["数据概览", "审核管理", "模型管理"]
srow1_annots = [
    "指标卡片 · 趋势图 · 分布图",
    "图片预览 · 批量操作 · 快捷键",
    "一键训练 · 类别详情 · 进度监控",
]
for i in range(3):
    x = (SLIDE_W - 3*S_IMG_W - 2*S_IMG_GAP)//2 + i*(S_IMG_W + S_IMG_GAP)
    add_placeholder(s4, x, SR1_Y, S_IMG_W, S_IMG_H)
    # 上方注释
    add_tb(s4, x, SR1_Y - Inches(0.42), S_IMG_W, Inches(0.4),
           srow1_annots[i], size=Pt(9), color=GRAY, align=PP_ALIGN.CENTER)
    add_tb(s4, x, SR1_Y - Inches(0.6), S_IMG_W, Inches(0.22),
           srow1_labels[i], size=Pt(10), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# 中行一个 — 系统配置
mid_x = (SLIDE_W - S_IMG_W) // 2
add_placeholder(s4, mid_x, SR2_Y, S_IMG_W, S_IMG_H)
add_tb(s4, mid_x - Inches(1.4), SR2_Y + Inches(0.3), Inches(1.3), Inches(0.7),
       "词条 CRUD · 配图上传\n拼音搜索",
       size=Pt(9), color=GRAY, align=PP_ALIGN.RIGHT)
add_tb(s4, mid_x + S_IMG_W + Inches(0.1), SR2_Y + Inches(0.3), Inches(1.3), Inches(0.7),
       "系统配置\nLLM API-Key · 连接测试\n上传限制",
       size=Pt(9), bold=True, color=BLUE)

# 下行三个
srow3_labels = ["百科管理", "用户管理", "系统日志"]
srow3_annots = [
    "词条 CRUD · 配图上传 · 拼音搜索",
    "用户列表 · 搜索 · 信息编辑",
    "日志查看 · 级别筛选 · 关键词",
]
for i in range(3):
    x = (SLIDE_W - 3*S_IMG_W - 2*S_IMG_GAP)//2 + i*(S_IMG_W + S_IMG_GAP)
    add_placeholder(s4, x, SR3_Y, S_IMG_W, S_IMG_H)
    # 下方注释
    add_tb(s4, x, SR3_Y + S_IMG_H + Inches(0.05), S_IMG_W, Inches(0.4),
           srow3_annots[i], size=Pt(9), color=GRAY, align=PP_ALIGN.CENTER)
    add_tb(s4, x, SR3_Y + S_IMG_H + Inches(0.38), S_IMG_W, Inches(0.22),
           srow3_labels[i], size=Pt(10), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

add_bottom_bar(s4, "管理员后台 — 密码登录 + Bearer Token 鉴权，7 个子页面覆盖全部管理场景")

# ═══════════════════════════════════════════════════════
# 第 5 页 — AI 协同开发
# ═══════════════════════════════════════════════════════
print("生成第 5 页...")
s5 = prs.slides.add_slide(TITLE_CONTENT_LAYOUT)
add_page_title(s5, 5, "AI 协同开发实践")

LEFT_W = Inches(4.6)
RIGHT_W = Inches(7.2)
RIGHT_X = Inches(5.4)
SECTION_H = Inches(1.85)
SECTION_GAP = Inches(0.15)
START_Y = Inches(1.05)

prompts = [
    ("① 架构设计阶段",
     "\"将 410 行单文件重构为路由/状态/API/\n 组件/页面五层分离，Vue3 + Pinia，\n 零构建 CDN 方案\"",
     "→ 产出 32 个文件的分层架构"),
    ("② 编码实现阶段",
     "\"设计农业科技风格 CSS 设计 Token：\n 玻璃拟态、3D 卡片、Mesh Gradient\n 极光背景，147 变量，5 级分层\"",
     "→ 产出 147 个设计 Token + 5 个 CSS 文件"),
    ("③ 联调完善阶段",
     "\"实现管理员后台 7 个子页面，\n 对接全部后端 API：仪表盘/审核/\n 模型管理/配置/百科/用户/日志\"",
     "→ 产出 1,900 行管理后台代码"),
]

for i, (phase_title, prompt_text, result_text) in enumerate(prompts):
    sy = START_Y + i * (SECTION_H + SECTION_GAP)

    # 阶段标题
    add_tb(s5, MARGIN, sy, LEFT_W, Inches(0.3),
           phase_title, size=Pt(16), bold=True, color=BLUE)

    # Prompt 内容（灰底）
    prompt_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      MARGIN, sy + Inches(0.32), LEFT_W, Inches(1.0))
    prompt_box.fill.solid()
    prompt_box.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    prompt_box.line.fill.background()
    add_tb(s5, MARGIN + Inches(0.15), sy + Inches(0.38), LEFT_W - Inches(0.3), Inches(0.9),
           prompt_text, size=Pt(12), color=DARK)

    # 产出结果
    add_tb(s5, MARGIN, sy + Inches(1.38), LEFT_W, Inches(0.35),
           result_text, size=Pt(13), bold=True, color=ACCENT_GREEN)

    # 分隔线
    if i < 2:
        line = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    MARGIN, sy + SECTION_H, LEFT_W, Pt(0.5))
        line.fill.solid()
        line.fill.fore_color.rgb = MLGRAY
        line.line.fill.background()

# 右侧三个 Claude Code 截图
SCREENSHOT_W = RIGHT_W - Inches(0.4)
SCREENSHOT_H = Inches(1.85)
for i in range(3):
    sy = START_Y + i * (SECTION_H + SECTION_GAP)
    add_placeholder(s5, RIGHT_X, sy, SCREENSHOT_W, SCREENSHOT_H,
                    f"Claude Code 对话截图 {i+1}")

# 底部协作模式
add_tb(s5, MARGIN, SLIDE_H - Inches(0.65), CONTENT_W, Inches(0.35),
       "协作模式：人 → 需求定义 + 代码审查 + 联调验证    |    AI → 架构生成 + 代码实现 + 样式产出",
       size=Pt(12), color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# 第 6 页 — 成果总结
# ═══════════════════════════════════════════════════════
print("生成第 6 页...")
s6 = prs.slides.add_slide(COVER_LAYOUT)

# 主标题
add_tb(s6, MARGIN, Inches(0.6), CONTENT_W, Inches(0.7),
       "成果总结", font=TITLE_FONT, size=Pt(48), color=DARK, align=PP_ALIGN.CENTER)

# 2×3 指标卡片
CARD_W6 = Inches(2.8)
CARD_H6 = Inches(1.0)
CARD_GAP6 = Inches(0.35)

items = [
    ("55 个", "源文件"), ("7,700 行", "代码"), ("13 个", "可交互页面"),
    ("6 大", "功能模块"), ("9 个", "通用组件"), ("147 个", "设计 Token"),
]

total_w = 3 * CARD_W6 + 2 * CARD_GAP6
start_x6 = (SLIDE_W - total_w) // 2
card_y = Inches(1.8)

for idx, (num, desc) in enumerate(items):
    row = idx // 3
    col = idx % 3
    x = start_x6 + col * (CARD_W6 + CARD_GAP6)
    y = card_y + row * (CARD_H6 + Inches(0.5))
    add_card(s6, x, y, CARD_W6, CARD_H6, num, desc)

# 中间说明文字
add_tb(s6, MARGIN, Inches(4.2), CONTENT_W, Inches(0.6),
       "桌面 / 平板 / 手机 — 三级响应式适配\n全部后端 API 已完成联调，系统可完整运行",
       size=Pt(15), color=GRAY, align=PP_ALIGN.CENTER)

# 技术栈
add_tb(s6, MARGIN, Inches(5.1), CONTENT_W, Inches(0.4),
       "Vue 3 · Pinia · ECharts  ·  零构建 · ESM importmap · CDN 直载  ·  Claude Code AI 协同",
       size=Pt(13), color=GRAY, align=PP_ALIGN.CENTER)

# 致谢
add_tb(s6, MARGIN, Inches(5.9), CONTENT_W, Inches(0.6),
       "感谢聆听 · 欢迎提问",
       font=TITLE_FONT, size=Pt(28), color=DARK, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# 保存
# ═══════════════════════════════════════════════════════
prs.save(SAVE_PATH)
print(f"\n✅ PPT 生成完成: {SAVE_PATH}")
print(f"   共 {len(prs.slides)} 页幻灯片")
os.startfile(os.path.dirname(SAVE_PATH))  # 打开文件夹
